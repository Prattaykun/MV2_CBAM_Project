
import sys
import os

# Install python-pptx if not installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        # Consistent Title styling
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear() # Clear default empty paragraph
        
        for text in content_text_list:
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(20)
            p.space_after = Pt(10)

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Implementation and Enhancement of MobileNetV2-CBAM for UAV-based Fire Detection"
    subtitle.text = "Author: Prattay Roy Chowdhury\nDate: February 10, 2026"

    # Slide 2: Problem Statement & Objective
    add_slide("Problem Statement & Objective", [
        "Problem: Forest fires are devastating, and traditional detection (satellites) is slow.",
        "Challenge: Lightweight UAV models struggle with false positives from smoke and fog.",
        "Objective: Build a robust, real-time fire detection system for edge devices.",
        "Key Goal: Solve the 'Generalization Gap' identified in literature."
    ])

    # Slide 3: Literature Review & Research Gap
    add_slide("Literature Review & Research Gap", [
        "Reference Paper: 'A lightweight CNN model for UAV-based image classification' (Soft Computing, 2025).",
        "Base Model: MobileNetV2 (Chosen for efficiency).",
        "Identified Gap (Page 14): 'Most error-prone fire images contain strong smoke.'",
        "Our Hypothesis: Integrating Attention Mechanisms (CBAM) + Domain Adaptation will fix this."
    ])

    # Slide 4: Proposed Methodology (Architecture)
    # We will draw a simple diagram here using shapes
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Proposed Architecture: MobileNetV2-CBAM"
    
    # Draw simple flow
    shapes = slide.shapes
    
    # Box 1: Input
    left = Inches(0.5); top = Inches(3.0); width = Inches(1.5); height = Inches(1.0)
    shape = shapes.add_shape(1, left, top, width, height) # 1 = RECTANGLE
    shape.text = "Input Image\n(224x224)"
    
    # Arrow 1
    shapes.add_shape(33, left + width, top + Inches(0.4), Inches(0.5), Inches(0.2)) # 33 = RIGHT_ARROW
    
    # Box 2: MobileNetV2
    left += Inches(2.0)
    shape = shapes.add_shape(1, left, top, width, height)
    shape.text = "MobileNetV2\n(Features)"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(100, 149, 237) # Cornflower Blue
    
    # Arrow 2
    shapes.add_shape(33, left + width, top + Inches(0.4), Inches(0.5), Inches(0.2))
    
    # Box 3: CBAM
    left += Inches(2.0)
    shape = shapes.add_shape(1, left, top, width, height)
    shape.text = "CBAM\n(Attention)"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 127, 80) # Coral
    
    # Arrow 3
    shapes.add_shape(33, left + width, top + Inches(0.4), Inches(0.5), Inches(0.2))
    
    # Box 4: Classifier
    left += Inches(2.0)
    shape = shapes.add_shape(1, left, top, width, height)
    shape.text = "Classifier\n(Fire / No Fire)"
     
    # Explanation Text
    txBox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "CBAM Mechanism: 1. Channel Attention (What?)  2. Spatial Attention (Where?)"

    # Slide 5: Technical Implementation
    add_slide("Technical Implementation", [
        "Frameworks: PyTorch (ML), FastAPI (Backend), Next.js (Dashboard).",
        "Hardware: Optimized for NVIDIA RTX 3050 (GPU Accelerated).",
        "Key Features:",
        "- Real-time Inference via API.",
        "- Grad-CAM Visualization for Explainability.",
        "- Simple Web Interface for UAV Operators."
    ])

    # Slide 6: Results & Comparative Analysis
    add_slide("Results & Comparative Analysis", [
        "Baseline (MobileNetV2): Failed on smoky Kaggle dataset (48.4% Accuracy).",
        "Proposed Solution (MobileNetV2-CBAM + DA): Achieved 95.8% Accuracy.",
        "Improvement: +47.4% boost in generalization.",
        "Conclusion: Domain Adaptation and Attention solve the smoke confusion issue."
    ])

    # Slide 7: Visualizations (Placeholder)
    slide_layout = prs.slide_layouts[5] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Visualizations (Demo)"
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    tf = txBox.text_frame
    tf.text = "[Insert Screenshot of Dashboard Here]\n[Insert Grad-CAM Heatmap Here]"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 8: Future Work & Recommendations
    add_slide("Future Work & Recommendations", [
        "1. Deployment: Port to NVIDIA Jetson Nano using TensorRT.",
        "2. Data: Collect Night-time and Foggy Fire images.",
        "3. Integration: Connect directly to drone flight controller (Pixhawk)."
    ])

    # Slide 9: Conclusion
    add_slide("Conclusion", [
        "We successfully enhanced MobileNetV2 for robust fire detection.",
        "Identified and fixed the critical vulnerability to smoke.",
        "Delivered a deployable, real-time system suitable for UAVs."
    ])
    
    # Slide 10: Suggested Journals
    add_slide("Suggested Journals (Q2/Q3)", [
         "1. Int. Journal of Machine Learning and Cybernetics (Springer) - Q2",
         "2. Journal of the Optical Society of America A - Q2",
         "3. Int. Journal of Computer Networks and Applications - Q3"
    ])

    output_path = os.path.join("docs", "Project_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
